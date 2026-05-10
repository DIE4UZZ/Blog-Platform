#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
毕业论文答辩PPT生成脚本
论文题目：融合内容推荐与数据分析的响应式博客平台设计与实现
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============ 配色方案 ============
PRIMARY_COLOR = RGBColor(0x1A, 0x56, 0xDB)      # 主色-深蓝
SECONDARY_COLOR = RGBColor(0x3B, 0x82, 0xF6)    # 辅色-蓝色
ACCENT_COLOR = RGBColor(0x10, 0xB9, 0x81)       # 强调色-绿色
DARK_COLOR = RGBColor(0x1F, 0x29, 0x37)         # 深色文字
LIGHT_COLOR = RGBColor(0xFF, 0xFF, 0xFF)        # 白色
GRAY_COLOR = RGBColor(0x6B, 0x72, 0x80)         # 灰色
BG_COLOR = RGBColor(0xF8, 0xFA, 0xFC)           # 浅灰背景
TITLE_BG = RGBColor(0x1E, 0x40, 0xAF)           # 标题背景深蓝


def set_slide_bg(slide, color):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_with_text(slide, left, top, width, height, text, 
                        font_size=12, font_color=DARK_COLOR, 
                        fill_color=None, bold=False, alignment=PP_ALIGN.LEFT,
                        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """添加带文字的形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Cm(0.5)
    tf.margin_right = Cm(0.5)
    tf.margin_top = Cm(0.3)
    tf.margin_bottom = Cm(0.3)
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    
    return shape


def add_title_bar(slide, title_text):
    """添加页面顶部标题栏"""
    # 标题背景条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        Inches(13.33), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # 标题文字
    tf = shape.text_frame
    tf.margin_left = Cm(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.color.rgb = LIGHT_COLOR
    p.font.bold = True


def add_bullet_points(slide, left, top, width, height, items, font_size=16, spacing=1.5):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(font_size * spacing)
        p.level = 0
    
    return txBox


def add_numbered_items(slide, left, top, width, height, items, font_size=15):
    """添加编号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(12)
    
    return txBox


def create_presentation():
    """创建答辩PPT"""
    prs = Presentation()
    # 设置16:9宽屏
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ========== 第1页：封面 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, PRIMARY_COLOR)
    
    # 装饰线条
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.8), Inches(2), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_COLOR
    line.line.fill.background()
    
    # 论文类型
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(10), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "本科毕业设计（论文）答辩"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    
    # 标题
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "融合内容推荐与数据分析的"
    p.font.size = Pt(36)
    p.font.color.rgb = LIGHT_COLOR
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "响应式博客平台设计与实现"
    p2.font.size = Pt(36)
    p2.font.color.rgb = LIGHT_COLOR
    p2.font.bold = True
    
    # 个人信息
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(10), Inches(1.8))
    tf = txBox.text_frame
    info_lines = [
        "答辩人：鄢子力",
        "指导教师：杨霁琳",
        "院    系：计算机科学学院 · 网络工程专业",
        "学    号：2022110456"
    ]
    for i, line_text in enumerate(info_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_text
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
        p.space_after = Pt(8)
    
    # ========== 第2页：目录 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "汇报提纲")
    
    outline_items = [
        ("01", "研究背景与意义", "信息过载问题与个性化推荐需求"),
        ("02", "相关技术介绍", "Vue3 + FastAPI + 混合推荐算法"),
        ("03", "系统设计", "架构设计、数据库设计、推荐流程设计"),
        ("04", "系统实现", "核心功能模块与关键代码实现"),
        ("05", "测试与效果分析", "功能测试、响应式适配、推荐效果评估"),
        ("06", "总结与展望", "工作总结与未来优化方向"),
    ]
    
    for i, (num, title, desc) in enumerate(outline_items):
        row = i // 2
        col = i % 2
        x = Inches(1.0 + col * 6.0)
        y = Inches(1.6 + row * 1.9)
        
        # 编号圆形
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY_COLOR
        circle.line.fill.background()
        ctf = circle.text_frame
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.text = num
        cp.font.size = Pt(18)
        cp.font.color.rgb = LIGHT_COLOR
        cp.font.bold = True
        cp.alignment = PP_ALIGN.CENTER
        
        # 标题和描述
        txBox = slide.shapes.add_textbox(x + Inches(1.0), y, Inches(4.5), Inches(0.9))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_COLOR
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY_COLOR
    
    # ========== 第3页：研究背景 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "一、研究背景与意义")
    
    # 问题描述
    problems = [
        "⚡ 信息过载：博客内容持续增长，用户筛选成本明显上升",
        "⚡ 传统展示方式局限：按时间/分类展示难以贴合用户兴趣",
        "⚡ 多设备适配需求：PC、平板、手机访问体验需统一",
        "⚡ 创作者缺乏数据支撑：无法了解内容传播效果与用户偏好",
    ]
    add_bullet_points(slide, Inches(0.8), Inches(1.5), Inches(6.5), Inches(3.5), problems, font_size=17)
    
    # 解决思路
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(12), Inches(2.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 解决思路"
    p.font.size = Pt(18)
    p.font.color.rgb = PRIMARY_COLOR
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = "设计一个融合内容推荐与数据分析的响应式博客平台，在保证系统轻量可实现的前提下，"
    p2.font.size = Pt(15)
    p2.font.color.rgb = DARK_COLOR
    p3 = tf.add_paragraph()
    p3.text = '提高内容分发的针对性和使用体验，形成"内容发布→用户反馈→推荐分发→效果分析→策略优化"的闭环。'
    p3.font.size = Pt(15)
    p3.font.color.rgb = DARK_COLOR
    
    # 右侧关键词
    keywords = ["个性化推荐", "行为分析", "响应式设计", "数据可视化", "混合推荐"]
    for i, kw in enumerate(keywords):
        shape = add_shape_with_text(
            slide, Inches(8.5), Inches(1.6 + i * 0.75), Inches(3.5), Inches(0.6),
            kw, font_size=14, font_color=LIGHT_COLOR, fill_color=SECONDARY_COLOR,
            bold=True, alignment=PP_ALIGN.CENTER
        )
    
    # ========== 第4页：技术栈 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "二、技术栈与架构概览")
    
    # 前端技术
    add_shape_with_text(slide, Inches(0.8), Inches(1.5), Inches(3.8), Inches(0.6),
                        "🖥️ 前端技术", font_size=18, font_color=PRIMARY_COLOR, bold=True)
    frontend_items = ["• Vue3 (Composition API)", "• Element-Plus 组件库", 
                      "• ECharts 数据可视化", "• Axios 网络请求", "• CSS3 Grid/Flex 响应式"]
    add_bullet_points(slide, Inches(0.8), Inches(2.2), Inches(3.8), Inches(3.0), frontend_items, font_size=14)
    
    # 后端技术
    add_shape_with_text(slide, Inches(4.8), Inches(1.5), Inches(3.8), Inches(0.6),
                        "⚙️ 后端技术", font_size=18, font_color=PRIMARY_COLOR, bold=True)
    backend_items = ["• Python + FastAPI", "• SQLAlchemy ORM", 
                     "• MySQL 8.0", "• JWT 认证鉴权", "• jieba 分词 + TF-IDF"]
    add_bullet_points(slide, Inches(4.8), Inches(2.2), Inches(3.8), Inches(3.0), backend_items, font_size=14)
    
    # 推荐算法
    add_shape_with_text(slide, Inches(8.8), Inches(1.5), Inches(4.0), Inches(0.6),
                        "🧠 推荐算法", font_size=18, font_color=PRIMARY_COLOR, bold=True)
    algo_items = ["• 基于内容推荐 (TF-IDF)", "• 协同过滤 (User-CF)", 
                  "• 热度信号融合", "• 混合加权排序", "• 冷启动推荐策略"]
    add_bullet_points(slide, Inches(8.8), Inches(2.2), Inches(4.0), Inches(3.0), algo_items, font_size=14)
    
    # 架构说明
    add_shape_with_text(slide, Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.5),
                        "系统架构：B/S架构 + 前后端分离 | 表现层(Vue3) → 业务逻辑层(FastAPI RESTful API) → 数据持久层(MySQL)",
                        font_size=15, font_color=DARK_COLOR, fill_color=RGBColor(0xE0, 0xE7, 0xFF),
                        alignment=PP_ALIGN.CENTER)
    
    # ========== 第5页：系统功能设计 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "三、系统功能设计")
    
    modules = [
        ("用户与认证", "注册、登录、JWT鉴权\n个人中心、角色管理", PRIMARY_COLOR),
        ("文章管理", "Markdown编辑发布\n分类标签、搜索检索", SECONDARY_COLOR),
        ("互动行为", "点赞、收藏、评论\n关注、通知", ACCENT_COLOR),
        ("混合推荐", "内容推荐+协同过滤\n热度融合+冷启动", RGBColor(0xF5, 0x9E, 0x0B)),
        ("数据分析", "阅读趋势、推荐效果\n用户画像、内容表现", RGBColor(0xEF, 0x44, 0x44)),
    ]
    
    for i, (title, desc, color) in enumerate(modules):
        x = Inches(0.5 + i * 2.5)
        y = Inches(1.8)
        
        # 模块卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.3), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_COLOR
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        # 模块标题
        title_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.3), Inches(2.1), Inches(0.6))
        ttf = title_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(16)
        tp.font.color.rgb = color
        tp.font.bold = True
        tp.alignment = PP_ALIGN.CENTER
        
        # 模块描述
        desc_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(1.0), Inches(2.1), Inches(1.3))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(12)
        dp.font.color.rgb = GRAY_COLOR
        dp.alignment = PP_ALIGN.CENTER
    
    # 底部说明
    add_shape_with_text(slide, Inches(0.8), Inches(5.0), Inches(11.8), Inches(2.0),
                        "五大模块协同工作：用户行为数据 → 推荐引擎计算 → 个性化内容分发 → 效果分析反馈 → 策略持续优化",
                        font_size=15, font_color=DARK_COLOR, fill_color=RGBColor(0xEC, 0xFD, 0xF5),
                        alignment=PP_ALIGN.CENTER)
    
    # ========== 第6页：数据库设计 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "四、数据库设计")
    
    # 核心表
    tables_info = [
        ("用户表 (user)", "id, username, email, password,\nrole, preference_tags, create_time"),
        ("文章表 (article)", "id, author_id, title, content,\ncategory, tags, view/like/collect_count"),
        ("行为表 (user_behavior)", "id, user_id, article_id,\nbehavior_type, read_duration, scroll_depth"),
        ("推荐表 (recommendation)", "id, user_id, article_id,\nrecommend_type, score, is_clicked"),
        ("评论表 (comment)", "id, article_id, user_id,\ncontent, parent_id"),
        ("点赞/收藏表", "id, user_id, article_id,\ncreate_time (唯一约束)"),
    ]
    
    for i, (tname, fields) in enumerate(tables_info):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.5 + row * 2.8)
        
        # 表名
        name_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(3.9), Inches(0.55))
        name_box.fill.solid()
        name_box.fill.fore_color.rgb = PRIMARY_COLOR
        name_box.line.fill.background()
        ntf = name_box.text_frame
        ntf.vertical_anchor = MSO_ANCHOR.MIDDLE
        np = ntf.paragraphs[0]
        np.text = tname
        np.font.size = Pt(13)
        np.font.color.rgb = LIGHT_COLOR
        np.font.bold = True
        np.alignment = PP_ALIGN.CENTER
        
        # 字段
        field_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(3.5), Inches(1.8))
        ftf = field_box.text_frame
        ftf.word_wrap = True
        fp = ftf.paragraphs[0]
        fp.text = fields
        fp.font.size = Pt(12)
        fp.font.color.rgb = GRAY_COLOR
    
    # 设计原则
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "设计原则：贴合业务需求 | 规范与扩展并重 | 优先查询性能 | 推荐与分析数据共用 | 安全与可维护"
    p.font.size = Pt(12)
    p.font.color.rgb = GRAY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # ========== 第7页：推荐算法设计 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "五、混合推荐算法设计")
    
    # 流程步骤
    steps = [
        ("1️⃣ 用户画像构建", "阅读时长→1~5分映射\n点赞=4分, 收藏/评论=5分\n+ 用户兴趣标签补充"),
        ("2️⃣ 内容召回", "文章文本拼接→jieba分词\n→TF-IDF向量化\n→余弦相似度计算"),
        ("3️⃣ 协同召回", "统计交互重叠用户\n→筛选相似用户\n→汇总偏好贡献"),
        ("4️⃣ 融合排序", "内容分×0.4 + 协同分×0.3\n+ 热度分×0.3\n(归一化后加权)"),
        ("5️⃣ 结果输出", "过滤已交互文章\n→Top-N分页返回\n→写入推荐记录表"),
    ]
    
    for i, (step_title, step_desc) in enumerate(steps):
        x = Inches(0.3 + i * 2.6)
        y = Inches(1.5)
        
        # 步骤卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.4), Inches(2.8))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_COLOR
        card.line.color.rgb = SECONDARY_COLOR
        card.line.width = Pt(1.5)
        
        # 步骤标题
        st_box = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.2), Inches(2.2), Inches(0.6))
        stf = st_box.text_frame
        sp = stf.paragraphs[0]
        sp.text = step_title
        sp.font.size = Pt(13)
        sp.font.color.rgb = PRIMARY_COLOR
        sp.font.bold = True
        sp.alignment = PP_ALIGN.CENTER
        
        # 步骤描述
        sd_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.9), Inches(2.1), Inches(1.8))
        sdf = sd_box.text_frame
        sdf.word_wrap = True
        sdp = sdf.paragraphs[0]
        sdp.text = step_desc
        sdp.font.size = Pt(11)
        sdp.font.color.rgb = DARK_COLOR
        sdp.alignment = PP_ALIGN.CENTER
    
    # 融合公式说明
    formula_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(12.3), Inches(2.5))
    ftf = formula_box.text_frame
    ftf.word_wrap = True
    
    fp = ftf.paragraphs[0]
    fp.text = "融合排序策略："
    fp.font.size = Pt(15)
    fp.font.color.rgb = PRIMARY_COLOR
    fp.font.bold = True
    
    formulas = [
        "• 热度分 = View×1 + Like×3 + Collect×5 + Comment×4",
        "• 内容+协同均存在：Final = 0.4×Content + 0.3×Collab + 0.3×Hot",
        "• 仅协同信号：Final = 0.6×Collab + 0.4×Hot",
        "• 仅内容信号：Final = 0.7×Content + 0.3×Hot",
        "• 冷启动退化：Final = Hot（按热度排序）",
    ]
    for f in formulas:
        p = ftf.add_paragraph()
        p.text = f
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(4)
    
    # ========== 第8页：行为采集实现 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "六、用户行为数据采集")
    
    # 左侧：采集类型
    add_shape_with_text(slide, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.5),
                        "📊 采集的行为类型", font_size=16, font_color=PRIMARY_COLOR, bold=True)
    
    behavior_items = [
        "显式行为：点赞、收藏、评论、搜索",
        "隐式行为：阅读时长、滚动深度",
        "",
        "📦 上报机制：",
        "• 本地队列缓存，阈值20条批量上报",
        "• 定时器8秒触发刷新",
        "• 失败重试（最多3次）",
        "• 页面关闭时 sendBeacon 兜底补传",
        "• 游客行为也记录（user_id可空）",
    ]
    add_bullet_points(slide, Inches(0.5), Inches(2.1), Inches(5.5), Inches(4.5), behavior_items, font_size=14)
    
    # 右侧：数据用途
    add_shape_with_text(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
                        "🎯 数据用途", font_size=16, font_color=PRIMARY_COLOR, bold=True)
    
    usage_items = [
        "→ 用户兴趣画像构建",
        "→ 推荐算法偏好建模",
        "→ 阅读趋势分析 (PV/UV)",
        "→ 用户画像分析",
        "→ 推荐效果评估 (CTR)",
        "→ 内容表现分析",
        "",
        "💡 设计理念：",
        "采集不影响用户体验",
        "数据完整性优先",
    ]
    add_bullet_points(slide, Inches(7.0), Inches(2.1), Inches(5.5), Inches(4.5), usage_items, font_size=14)
    
    # ========== 第9页：数据分析看板 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "七、数据可视化分析看板")
    
    # 四个分析模块
    analysis_modules = [
        ("📈 阅读趋势", "PV、UV、平均阅读时长\n平均滚动深度\n支持日/周/月粒度", Inches(0.5), Inches(1.5)),
        ("🎯 推荐效果", "曝光量、点击量、CTR\n转化率、来源分布对比\n趋势变化分析", Inches(6.8), Inches(1.5)),
        ("👤 用户画像", "兴趣标签、行为分布\n活跃时段热力图\n偏好分类、搜索词", Inches(0.5), Inches(4.2)),
        ("📊 内容表现", "文章阅读量排行\n互动率统计\n分类热度雷达图", Inches(6.8), Inches(4.2)),
    ]
    
    for title, desc, x, y in analysis_modules:
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(5.8), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_COLOR
        card.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        card.line.width = Pt(1)
        
        # 标题
        t_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.6))
        ttf = t_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(18)
        tp.font.color.rgb = PRIMARY_COLOR
        tp.font.bold = True
        
        # 描述
        d_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), Inches(5.2), Inches(1.3))
        dtf = d_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(14)
        dp.font.color.rgb = DARK_COLOR
    
    # ========== 第10页：响应式设计 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "八、响应式设计与多端适配")
    
    # 左侧说明
    responsive_items = [
        "🖥️ PC端（≥1024px）：",
        "   多栏布局，侧边栏+主内容区",
        "   完整导航栏，图表并排展示",
        "",
        "📱 移动端（≤640px）：",
        "   单栏布局，汉堡菜单导航",
        "   图表自适应高度，触控优化",
        "",
        "🔧 实现方式：",
        "   CSS3 Grid + Flex 弹性布局",
        "   @media 断点适配",
        "   repeat(auto-fit, minmax(...)) 自动排版",
        "   Element-Plus 响应式组件",
    ]
    add_bullet_points(slide, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5), responsive_items, font_size=14)
    
    # 右侧效果说明
    add_shape_with_text(slide, Inches(7.5), Inches(1.8), Inches(5.0), Inches(4.5),
                        "✅ 适配效果\n\n"
                        "• PC 1920×1080 正常展示\n"
                        "• 移动端 375×812 正常展示\n"
                        "• 无功能失效或界面混乱\n"
                        "• 文字/图片自适应缩放\n"
                        "• 导航模式自动切换\n\n"
                        "测试结果：跨屏兼容性良好",
                        font_size=14, font_color=DARK_COLOR, 
                        fill_color=RGBColor(0xF0, 0xFD, 0xF4))
    
    # ========== 第11页：测试结果 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "九、系统测试与推荐效果评估")
    
    # 功能测试
    add_shape_with_text(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
                        "✅ 功能测试结果", font_size=16, font_color=ACCENT_COLOR, bold=True)
    
    test_items = [
        "• 注册/登录流程正常，JWT鉴权有效",
        "• 文章发布、编辑、删除功能完整",
        "• 点赞/收藏/评论实时更新",
        "• 行为采集日志完整记录",
        "• 推荐列表正常生成与展示",
        "• 数据分析看板指标准确",
    ]
    add_bullet_points(slide, Inches(0.5), Inches(2.1), Inches(6.0), Inches(3.5), test_items, font_size=14)
    
    # 推荐效果评估
    add_shape_with_text(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5),
                        "📊 推荐算法离线评估", font_size=16, font_color=PRIMARY_COLOR, bold=True)
    
    # 评估指标卡片
    metrics = [
        ("Precision@10", "0.1385", "推荐准确率"),
        ("Recall@10", "0.2279", "推荐召回率"),
        ("HitRate@10", "0.7692", "命中率"),
    ]
    
    for i, (metric_name, value, desc) in enumerate(metrics):
        y = Inches(2.2 + i * 1.4)
        
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                       Inches(7.2), y, Inches(5.3), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_COLOR
        card.line.color.rgb = SECONDARY_COLOR
        card.line.width = Pt(1)
        
        # 指标名
        m_box = slide.shapes.add_textbox(Inches(7.4), y + Inches(0.1), Inches(2.5), Inches(0.5))
        mtf = m_box.text_frame
        mp = mtf.paragraphs[0]
        mp.text = metric_name
        mp.font.size = Pt(13)
        mp.font.color.rgb = GRAY_COLOR
        
        # 数值
        v_box = slide.shapes.add_textbox(Inches(7.4), y + Inches(0.5), Inches(2.5), Inches(0.6))
        vtf = v_box.text_frame
        vp = vtf.paragraphs[0]
        vp.text = value
        vp.font.size = Pt(24)
        vp.font.color.rgb = PRIMARY_COLOR
        vp.font.bold = True
        
        # 描述
        d_box = slide.shapes.add_textbox(Inches(10.0), y + Inches(0.3), Inches(2.3), Inches(0.6))
        dtf = d_box.text_frame
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(13)
        dp.font.color.rgb = GRAY_COLOR
    
    # 评估说明
    eval_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2))
    etf = eval_box.text_frame
    etf.word_wrap = True
    ep = etf.paragraphs[0]
    ep.text = "评估条件：Top-K=10, 协同近邻数=20, 内容权重0.6/协同权重0.4, 用户最小交互数≥5"
    ep.font.size = Pt(12)
    ep.font.color.rgb = GRAY_COLOR
    ep2 = etf.add_paragraph()
    ep2.text = "结论：HitRate@10达76.92%，说明混合推荐方案在中小型博客平台场景下具备实际可用性"
    ep2.font.size = Pt(14)
    ep2.font.color.rgb = DARK_COLOR
    ep2.font.bold = True
    
    # ========== 第12页：创新点与贡献 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "十、主要工作与创新点")
    
    innovations = [
        ("🔹 完整业务闭环", 
         '实现了"内容发布→行为采集→推荐分发→效果分析→策略优化"的完整链路，\n而非仅实现单一推荐功能'),
        ("🔹 轻量混合推荐方案", 
         "结合TF-IDF内容相似度+协同过滤+热度信号，无需大模型和高算力，\n适合中小规模博客场景落地"),
        ("🔹 多维行为采集", 
         "不仅采集显式反馈，还记录阅读时长、滚动深度等隐式行为，\n配合批量上报+兜底补传机制保证数据完整性"),
        ("🔹 可视化分析看板", 
         "为创作者提供阅读趋势、推荐效果、用户画像、内容表现等多维分析，\n支撑数据驱动的内容运营决策"),
        ("🔹 响应式全端适配", 
         "基于CSS3 Grid/Flex实现PC与移动端自适应布局，\n保证不同设备上的一致使用体验"),
    ]
    
    for i, (title, desc) in enumerate(innovations):
        y = Inches(1.4 + i * 1.15)
        
        t_box = slide.shapes.add_textbox(Inches(0.8), y, Inches(11.5), Inches(0.4))
        ttf = t_box.text_frame
        tp = ttf.paragraphs[0]
        tp.text = title
        tp.font.size = Pt(16)
        tp.font.color.rgb = PRIMARY_COLOR
        tp.font.bold = True
        
        d_box = slide.shapes.add_textbox(Inches(1.2), y + Inches(0.4), Inches(11.0), Inches(0.7))
        dtf = d_box.text_frame
        dtf.word_wrap = True
        dp = dtf.paragraphs[0]
        dp.text = desc
        dp.font.size = Pt(13)
        dp.font.color.rgb = DARK_COLOR
    
    # ========== 第13页：总结与展望 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_COLOR)
    add_title_bar(slide, "十一、总结与展望")
    
    # 总结
    add_shape_with_text(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(0.5),
                        "📋 工作总结", font_size=17, font_color=PRIMARY_COLOR, bold=True)
    
    summary_items = [
        "✓ 完成前后端分离的响应式博客平台",
        "✓ 实现混合推荐引擎（内容+协同+热度）",
        "✓ 建立多维用户行为采集体系",
        "✓ 构建数据可视化分析看板",
        "✓ 离线评估验证推荐方案可用性",
        "✓ 形成完整的推荐-分析闭环",
    ]
    add_bullet_points(slide, Inches(0.5), Inches(2.1), Inches(6.0), Inches(3.5), summary_items, font_size=14)
    
    # 展望
    add_shape_with_text(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(0.5),
                        "🚀 未来展望", font_size=17, font_color=ACCENT_COLOR, bold=True)
    
    future_items = [
        "→ 引入BERT等预训练模型提升语义理解",
        "→ 建立在线A/B测试评估体系",
        "→ 引入Redis缓存提升推荐响应速度",
        "→ 补充匿名设备标识覆盖游客行为",
        "→ 索引优化→缓存分流→读写分离",
    ]
    add_bullet_points(slide, Inches(7.0), Inches(2.1), Inches(5.8), Inches(3.5), future_items, font_size=14)
    
    # 底部总结语
    add_shape_with_text(slide, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2),
                        "本文工作表明：在不依赖高成本模型和大规模算力的条件下，结合内容特征、用户行为与可视化分析，\n"
                        "仍然可以实现一套可运行、可分析、可迭代的个性化博客推荐系统。",
                        font_size=15, font_color=DARK_COLOR, fill_color=RGBColor(0xEF, 0xF6, 0xFF),
                        alignment=PP_ALIGN.CENTER)
    
    # ========== 第14页：致谢 ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY_COLOR)
    
    # 感谢标题
    txBox = slide.shapes.add_textbox(Inches(0), Inches(2.5), Inches(13.33), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢各位老师的指导与聆听！"
    p.font.size = Pt(36)
    p.font.color.rgb = LIGHT_COLOR
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    p2 = tf.add_paragraph()
    p2.text = ""
    p3 = tf.add_paragraph()
    p3.text = "恳请各位老师批评指正"
    p3.font.size = Pt(22)
    p3.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p3.alignment = PP_ALIGN.CENTER
    
    # 个人信息
    txBox2 = slide.shapes.add_textbox(Inches(0), Inches(5.0), Inches(13.33), Inches(1.5))
    tf2 = txBox2.text_frame
    p4 = tf2.paragraphs[0]
    p4.text = "答辩人：鄢子力 | 指导教师：杨霁琳"
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
    p4.alignment = PP_ALIGN.CENTER
    p5 = tf2.add_paragraph()
    p5.text = "计算机科学学院 · 网络工程专业 · 2022110456"
    p5.font.size = Pt(14)
    p5.font.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
    p5.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    output_path = "/Users/mac/Desktop/毕业论文答辩PPT_鄢子力.pptx"
    prs.save(output_path)
    print(f"✅ 答辩PPT已生成：{output_path}")
    print(f"   共 {len(prs.slides)} 页幻灯片")
    return output_path


if __name__ == "__main__":
    create_presentation()
